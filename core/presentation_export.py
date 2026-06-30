"""
Subject Analytics — Presentation Export Module
Generates downloadable PowerPoint (.pptx) and PDF (.pdf) reports
from the subject analytics context dict.

Dependencies:
  - python-pptx  (for .pptx generation)
  - reportlab     (for .pdf generation, already in project)
"""

import io
import logging
from datetime import date

from django.http import HttpResponse

logger = logging.getLogger(__name__)

# ─────────────────────────────────────────────────────────────────────────────
# Colour palette (matches web charts and AcadStat design system)
# ─────────────────────────────────────────────────────────────────────────────
COLOR_PRIMARY      = "#4F46E5"   # Indigo accent
COLOR_SUCCESS      = "#10B981"   # Green
COLOR_DANGER       = "#EF4444"   # Red
COLOR_WARNING      = "#F59E0B"   # Amber
COLOR_INFO         = "#3B82F6"   # Blue
COLOR_EXCELLENCE   = "#6366F1"   # Indigo-light
COLOR_TEXT_PRIMARY = "#1E293B"
COLOR_TEXT_SECONDARY = "#64748B"
COLOR_TEXT_MUTED   = "#94A3B8"
COLOR_WHITE        = "#FFFFFF"
COLOR_HEADER_BG    = "#EEF2FF"
COLOR_CHART_BG     = [
    "#667eea", "#10b981", "#f59e0b", "#ef4444",
    "#8b5cf6", "#06b6d4", "#ec4899", "#84cc16",
]

# ─────────────────────────────────────────────────────────────────────────────
# Shared helpers
# ─────────────────────────────────────────────────────────────────────────────

def _filter_label(context):
    """Build a short filter-description string for slide/page footers."""
    parts = []
    subj = context.get("subject_obj")
    parts.append(f"Subject={subj.name if subj else 'All'}")
    parts.append(f"Terminal={context.get('terminal', 'all')}")
    parts.append(f"Mode={'Single' if context.get('analysis_mode') == 'single' else 'Bulk'}")
    cs = context.get("class_section", "")
    parts.append(f"Class={cs or 'All'}")
    if context.get("selected_student"):
        parts.append(f"Student={context['selected_student'].name}")
    return " | ".join(parts)


def _filter_rows(context):
    """Return readable filter metadata rows for generated reports."""
    subj = context.get("subject_obj")
    selected_student = context.get("selected_student")
    class_section = context.get("class_section", "").replace("|", " - ") or "All sections"
    return [
        ("Analysis Mode", "Single Analysis" if context.get("analysis_mode") == "single" else "Bulk Analysis"),
        ("Subject", subj.name if subj else "All subjects assigned to teacher"),
        ("Terminal", context.get("terminal", "all")),
        ("Class / Section", class_section),
        ("Student", selected_student.name if selected_student else "All students in selected scope"),
    ]


def _scope_sentence(context, teacher_name):
    """Describe the exact teacher/filter scope in one sentence."""
    rows = dict(_filter_rows(context))
    return (
        f"Generated for {teacher_name}. Scope: {rows['Analysis Mode']}; "
        f"{rows['Subject']}; terminal {rows['Terminal']}; "
        f"{rows['Class / Section']}; {rows['Student']}."
    )


def _filename(context, ext):
    """Generate a dynamic filename e.g. SubjectAnalytics_Class10_Mathematics_2nd_2026-06-30.pptx"""
    cs = context.get("class_section", "all").replace(" ", "_").replace("|", "-") or "all"
    subj = context.get("subject_obj")
    subj_name = subj.name.replace(" ", "_") if subj else "all"
    term = context.get("terminal", "all")
    today = date.today().isoformat()
    return f"SubjectAnalytics_{cs}_{subj_name}_{term}_{today}.{ext}"


def _pct_color(pct):
    """Return hex colour for a percentage value."""
    if pct is None:
        return COLOR_TEXT_MUTED
    if pct >= 80:
        return COLOR_EXCELLENCE
    if pct >= 60:
        return COLOR_SUCCESS
    if pct >= 40:
        return COLOR_WARNING
    return COLOR_DANGER


def _teacher_display_name(teacher):
    """Return the best readable teacher name for exports."""
    if hasattr(teacher, "get_full_name"):
        name = teacher.get_full_name()
        if name:
            return name
    name = getattr(teacher, "name", "")
    if name:
        return name
    return str(teacher) if teacher else "Teacher"


# ═════════════════════════════════════════════════════════════════════════════
#  PPTX EXPORT
# ═════════════════════════════════════════════════════════════════════════════

def export_subject_analytics_pptx(request, context):
    """
    Build a PowerPoint presentation from the subject analytics context.
    Returns an HttpResponse with content-type
    'application/vnd.openxmlformats-officedocument.presentationml.presentation'.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.dml import MSO_THEME_COLOR
    except ImportError:
        logger.exception("python-pptx is not installed")
        return HttpResponse(
            "PowerPoint generation requires python-pptx. "
            "Please run: pip install python-pptx",
            status=500,
        )

    def _rgb(hex_color):
        """Convert hex string '#4F46E5' to RGBColor."""
        h = hex_color.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout

    # ── helpers ──────────────────────────────────────────────────────────
    def add_title(slide, text, left=0.5, top=0.3, width=12, height=0.8):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = _rgb(COLOR_TEXT_PRIMARY)
        return tf

    def add_subtitle(slide, text, left=0.5, top=1.1, width=12, height=0.5):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = _rgb(COLOR_TEXT_SECONDARY)
        return tf

    def add_footer(slide, text, left=0.5, top=7.0, width=12, height=0.4):
        line = slide.shapes.add_shape(
            1,
            Inches(0.5),
            Inches(6.88),
            Inches(12.3),
            Inches(0.02),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = _rgb(COLOR_HEADER_BG)
        line.line.fill.background()
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(9)
        p.font.color.rgb = _rgb(COLOR_TEXT_MUTED)
        return tf

    def add_slide_frame(slide):
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb("#F8FAFC")
        bg.line.fill.background()
        side = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
        side.fill.solid()
        side.fill.fore_color.rgb = _rgb(COLOR_PRIMARY)
        side.line.fill.background()
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.14))
        header.fill.solid()
        header.fill.fore_color.rgb = _rgb(COLOR_PRIMARY)
        header.line.fill.background()
        accent = slide.shapes.add_shape(1, Inches(0), Inches(0.14), Inches(13.333), Inches(0.05))
        accent.fill.solid()
        accent.fill.fore_color.rgb = _rgb(COLOR_SUCCESS)
        accent.line.fill.background()

    def add_panel(slide, left, top, width, height, fill="#FFFFFF", line=COLOR_HEADER_BG):
        panel = slide.shapes.add_shape(
            1, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = _rgb(fill)
        panel.line.color.rgb = _rgb(line)
        panel.line.width = Pt(1)
        return panel

    def add_section_label(slide, text, left, top, width=3.4, color=COLOR_PRIMARY):
        label = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.32))
        label.fill.solid()
        label.fill.fore_color.rgb = _rgb(color)
        label.line.fill.background()
        tf = label.text_frame
        tf.margin_left = Inches(0.08)
        p = tf.paragraphs[0]
        p.text = text.upper()
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = _rgb(COLOR_WHITE)
        return label

    def add_filter_strip(slide, rows, left=0.5, top=1.55, width=12.3):
        chip_width = width / 3
        for idx, (label, value) in enumerate(rows):
            row = idx // 3
            col = idx % 3
            x = left + (col * chip_width)
            y = top + (row * 0.48)
            box = slide.shapes.add_shape(
                1,
                Inches(x),
                Inches(y),
                Inches(chip_width - 0.12),
                Inches(0.34),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = _rgb("#F8FAFC")
            box.line.color.rgb = _rgb(COLOR_HEADER_BG)
            tf = box.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            p = tf.paragraphs[0]
            p.text = f"{label}: {value}"
            p.font.size = Pt(8.5)
            p.font.color.rgb = _rgb(COLOR_TEXT_SECONDARY)
            p.font.bold = True

    def add_metric_cards(slide, cards, left=0.5, top=2.3, width=12.3):
        card_width = width / max(len(cards), 1)
        for idx, (label, value, color) in enumerate(cards):
            x = left + idx * card_width
            shape = slide.shapes.add_shape(
                1,
                Inches(x),
                Inches(top),
                Inches(card_width - 0.12),
                Inches(1.05),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb("#FFFFFF")
            shape.line.color.rgb = _rgb(color)
            shape.line.width = Pt(1.4)
            tf = shape.text_frame
            tf.margin_left = Inches(0.12)
            tf.margin_right = Inches(0.12)
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = _rgb(COLOR_TEXT_SECONDARY)
            p2 = tf.add_paragraph()
            p2.text = str(value)
            p2.font.size = Pt(22)
            p2.font.bold = True
            p2.font.color.rgb = _rgb(color)

    def add_body_textbox(slide, left, top, width, height):
        return slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )

    def set_paragraph(tf, text, size=14, bold=False, color=COLOR_TEXT_PRIMARY, spacing_after=6):
        if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = _rgb(color)
        p.space_after = Pt(spacing_after)
        return p

    def add_table(slide, data, left=0.5, top=1.6, width=12, height=None):
        """Add a table to a slide. data[0] = header row."""
        rows = len(data)
        cols = len(data[0]) if data else 1
        if height is None:
            height = min(rows * 0.45, 5.0)
        table_shape = slide.shapes.add_table(
            rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        table = table_shape.table

        for r_idx, row_data in enumerate(data):
            for c_idx, cell_text in enumerate(row_data):
                cell = table.cell(r_idx, c_idx)
                cell.text = str(cell_text) if cell_text is not None else "—"
                # Style each paragraph in the cell
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(11)
                    if r_idx == 0:
                        para.font.bold = True
                        para.font.color.rgb = _rgb(COLOR_WHITE)
                        para.font.size = Pt(12)
                    else:
                        para.font.color.rgb = _rgb(COLOR_TEXT_PRIMARY)
                    para.alignment = PP_ALIGN.CENTER

                # Cell background
                if r_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(COLOR_PRIMARY)
                elif r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb("#F8FAFC")

        return table

    def new_slide(title_text, subtitle_text=None, footer_text=None):
        slide = prs.slides.add_slide(blank_layout)
        add_slide_frame(slide)
        add_title(slide, title_text)
        if subtitle_text:
            add_subtitle(slide, subtitle_text)
        if footer_text:
            add_footer(slide, footer_text)
        return slide

    # ── Extract context ──────────────────────────────────────────────────
    ctx = context
    teacher_name = _teacher_display_name(getattr(request, "teacher", None))
    filter_str = _filter_label(ctx)
    filter_rows = _filter_rows(ctx)
    today_str = date.today().strftime("%B %d, %Y")
    subj_obj = ctx.get("subject_obj")
    subj_display = subj_obj.name if subj_obj else "All Subjects"
    term_display = ctx.get("terminal", "all")
    mode_display = "Single Analysis" if ctx.get("analysis_mode") == "single" else "Bulk Analysis"
    cs_display = ctx.get("class_section", "").replace("|", " — ") or "All Sections"
    selected_student = ctx.get("selected_student")
    student_display = selected_student.name if selected_student else ""
    graph_summary = ctx.get("teacher_graph_summary", {})

    total = ctx.get("total", 0)

    # ══════════════════════════════════════════════════════════════════════════
    #  SLIDE 1: Title
    # ══════════════════════════════════════════════════════════════════════════
    slide = new_slide(
        "Subject Analytics Report",
        subtitle_text=f"{cs_display} | {subj_display} | {term_display} | {mode_display}",
        footer_text=f"Generated by AcadStat — {today_str} | Teacher: {teacher_name}",
    )
    add_filter_strip(slide, filter_rows, top=1.75)
    # Add a larger teacher/date block
    txBox = add_body_textbox(slide, 0.5, 3.0, 12, 1.2)
    tf = txBox.text_frame
    tf.word_wrap = True
    set_paragraph(tf, f"Teacher: {teacher_name}", size=18, bold=True, color=COLOR_TEXT_PRIMARY)
    set_paragraph(tf, f"Generated: {today_str}", size=13, color=COLOR_TEXT_SECONDARY)
    set_paragraph(
        tf,
        _scope_sentence(ctx, teacher_name),
        size=12,
        color=COLOR_TEXT_SECONDARY,
    )
    add_metric_cards(
        slide,
        [
            ("Average", f"{ctx.get('avg_pct', 0)}%", COLOR_PRIMARY),
            ("Pass Rate", f"{ctx.get('pass_rate', 0)}%", COLOR_SUCCESS),
            ("Highest", f"{ctx.get('highest_pct', 0)}%", COLOR_INFO),
            ("Below 40%", str(ctx.get("below_40", 0)), COLOR_DANGER),
        ],
        top=4.55,
    )

    # If no data, show a single-slide report
    if total == 0:
        slide2 = new_slide(
            "No Data Available",
            subtitle_text="The selected filters returned no results.",
            footer_text=filter_str,
        )
        txBox = add_body_textbox(slide2, 0.5, 2.5, 12, 2)
        tf = txBox.text_frame
        tf.word_wrap = True
        set_paragraph(
            tf,
            "No marks data was found for the current filter combination. "
            "Try adjusting your filters (subject, terminal, class section) or "
            "add marks from the 'Add Marks' page.",
            size=16,
            color=COLOR_TEXT_SECONDARY,
        )
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        resp = HttpResponse(
            buf.read(),
            content_type=(
                "application/vnd.openxmlformats-officedocument."
                "presentationml.presentation"
            ),
        )
        resp["Content-Disposition"] = (
            f'attachment; filename="{_filename(ctx, "pptx")}"'
        )
        return resp

    # ══════════════════════════════════════════════════════════════════════════
    #  SLIDE 2: Executive Summary (KPI Dashboard)
    # ══════════════════════════════════════════════════════════════════════════
    avg_pct = ctx.get("avg_pct", 0)
    highest_pct = ctx.get("highest_pct", 0)
    lowest_pct = ctx.get("lowest_pct", 0)
    pass_rate = ctx.get("pass_rate", 0)
    pass_count = ctx.get("pass_count", 0)
    below_40 = ctx.get("below_40", 0)

    slide = new_slide("Executive Summary - Key Metrics", footer_text=filter_str)
    add_filter_strip(slide, filter_rows, top=1.18)
    add_metric_cards(
        slide,
        [
            ("Average Score", f"{avg_pct}%", COLOR_PRIMARY),
            ("Pass Rate", f"{pass_rate}%", COLOR_SUCCESS),
            ("Highest Score", f"{highest_pct}%", COLOR_INFO),
            ("Lowest Score", f"{lowest_pct}%", COLOR_WARNING),
            ("Below 40%", str(below_40), COLOR_DANGER),
        ],
        top=2.15,
    )
    kpi_data = [
        ["Metric", "Value", "Metric", "Value"],
        [
            "Average Score",
            f"{avg_pct}%",
            "Pass Rate",
            f"{pass_rate}% ({pass_count}/{total})",
        ],
        [
            "Highest Score",
            f"{highest_pct}%",
            "Lowest Score",
            f"{lowest_pct}%",
        ],
        [
            "Below 40%",
            f"{below_40} results",
            "Analysis Mode",
            mode_display,
        ],
    ]
    if selected_student:
        kpi_data.append(["Student", student_display, "", ""])
    add_table(slide, kpi_data, top=3.65, height=2.2)

    # ══════════════════════════════════════════════════════════════════════════
    #  SLIDE 3: Score Distribution
    # ══════════════════════════════════════════════════════════════════════════
    dist_chart = ctx.get("distribution_chart", {})
    dist_labels = dist_chart.get("labels", [])
    dist_data = dist_chart.get("datasets", [{}])[0].get("data", [])
    dist_bg = dist_chart.get("datasets", [{}])[0].get("backgroundColor", [])

    slide = new_slide("Score Distribution", footer_text=filter_str)
    add_subtitle(
        slide,
        "Distribution of student scores across grade bands",
        top=1.0,
        height=0.4,
    )
    header_row = ["Grade Band", "Count", "Percentage of Total"]
    table_data = [header_row]
    for i, label in enumerate(dist_labels):
        count = dist_data[i] if i < len(dist_data) else 0
        pct_of_total = round((count / total * 100), 1) if total else 0
        table_data.append([label, str(count), f"{pct_of_total}%"])
    tbl = add_table(slide, table_data, top=1.6, height=3.0)

    # Colour-code the grade band cells
    if len(table_data) > 1:
        for r_idx in range(1, len(table_data)):
            cell = tbl.cell(r_idx, 0)
            bg_color = (
                dist_bg[r_idx - 1] if r_idx - 1 < len(dist_bg) else COLOR_PRIMARY
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(bg_color)
            for para in cell.text_frame.paragraphs:
                para.font.color.rgb = _rgb(COLOR_WHITE)
                para.font.bold = True

    add_panel(slide, 0.5, 4.95, 12, 1.45)
    add_section_label(slide, "Interpretation", 0.7, 5.12, color=COLOR_INFO)
    txBox = add_body_textbox(slide, 0.75, 5.48, 11.4, 0.7)
    tf = txBox.text_frame
    tf.word_wrap = True
    reading = graph_summary.get("graph_reading", [])
    dist_note = reading[1] if len(reading) > 1 else (
        f"{below_40} result(s) are below 40%, while "
        f"{dist_data[3] if len(dist_data) > 3 else 0} reached the 80-100% band."
    )
    set_paragraph(tf, dist_note, size=12, color=COLOR_TEXT_PRIMARY)

    # ══════════════════════════════════════════════════════════════════════════
    #  SLIDE 4: Terminal-Wise Performance
    # ══════════════════════════════════════════════════════════════════════════
    terminal_wise = ctx.get("terminal_wise_chart", {})
    term_labels = terminal_wise.get("labels", [])
    term_datasets = terminal_wise.get("datasets", [])

    slide = new_slide("Terminal-Wise Performance Trend", footer_text=filter_str)
    add_subtitle(
        slide,
        "Average percentage per terminal by subject",
        top=1.0,
        height=0.4,
    )
    if term_datasets:
        # Build table: rows = subjects, cols = terminals
        term_table_header = ["Subject"] + list(term_labels)
        term_table_data = [term_table_header]
        for ds in term_datasets:
            row = [ds.get("label", "Unknown")]
            vals = ds.get("data", [])
            for v in vals:
                row.append(f"{v}%" if v is not None else "—")
            term_table_data.append(row)
        tbl = add_table(slide, term_table_data, top=1.6, height=min(len(term_table_data) * 0.45, 4.5))

        # Colour-code percentage cells
        if len(term_table_data) > 1:
            for r_idx in range(1, len(term_table_data)):
                for c_idx in range(1, len(term_table_data[r_idx])):
                    cell = tbl.cell(r_idx, c_idx)
                    val = term_datasets[r_idx - 1].get("data", [])[c_idx - 1] if c_idx - 1 < len(term_datasets[r_idx - 1].get("data", [])) else None
                    if val is not None:
                        bg = _pct_color(val)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _rgb(bg)
                        for para in cell.text_frame.paragraphs:
                            para.font.color.rgb = _rgb(COLOR_WHITE)
                            para.font.bold = True
    else:
        txBox = add_body_textbox(slide, 0.5, 2.0, 12, 1)
        tf = txBox.text_frame
        set_paragraph(tf, "No terminal performance data available for current filters.", size=14, color=COLOR_TEXT_SECONDARY)

    add_panel(slide, 0.5, 6.05, 12, 0.75)
    txBox = add_body_textbox(slide, 0.7, 6.18, 11.6, 0.42)
    tf = txBox.text_frame
    tf.word_wrap = True
    trend_note = graph_summary.get("graph_reading", [""])[-1] if graph_summary.get("graph_reading") else "Terminal rows reflect the same logged-in teacher scope and current filters."
    set_paragraph(tf, trend_note, size=10.5, color=COLOR_TEXT_SECONDARY)

    # ══════════════════════════════════════════════════════════════════════════
    #  SLIDE 5: Subject Breakdown
    # ══════════════════════════════════════════════════════════════════════════
    subj_breakdown = ctx.get("subject_breakdown_chart", {})
    subj_labels = subj_breakdown.get("labels", [])
    subj_values = subj_breakdown.get("datasets", [{}])[0].get("data", [])
    subj_bg = subj_breakdown.get("datasets", [{}])[0].get("backgroundColor", COLOR_CHART_BG)

    slide = new_slide("Subject Breakdown", footer_text=filter_str)
    add_subtitle(
        slide,
        "Average percentage by subject in latest/selected terminal",
        top=1.0,
        height=0.4,
    )
    subj_header = ["Subject", "Average %"]
    subj_table_data = [subj_header]
    for i, label in enumerate(subj_labels):
        val = subj_values[i] if i < len(subj_values) else None
        subj_table_data.append([label, f"{val}%" if val is not None else "—"])
    tbl = add_table(slide, subj_table_data, top=1.6, width=6.1, height=min(len(subj_table_data) * 0.45, 4.0))

    # Colour-code subject cells
    if len(subj_table_data) > 1:
        for r_idx in range(1, len(subj_table_data)):
            cell = tbl.cell(r_idx, 0)
            if r_idx - 1 < len(subj_bg):
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(subj_bg[r_idx - 1])
                for para in cell.text_frame.paragraphs:
                    para.font.color.rgb = _rgb(COLOR_WHITE)
                    para.font.bold = True
            cell2 = tbl.cell(r_idx, 1)
            val = subj_values[r_idx - 1] if r_idx - 1 < len(subj_values) else None
            if val is not None:
                cell2.fill.solid()
                cell2.fill.fore_color.rgb = _rgb(_pct_color(val))
                for para in cell2.text_frame.paragraphs:
                    para.font.color.rgb = _rgb(COLOR_WHITE)
                    para.font.bold = True

    add_panel(slide, 6.9, 1.6, 5.6, 4.0, fill="#FFFFFF")
    add_section_label(slide, "Subject Focus", 7.12, 1.82, color=COLOR_SUCCESS)
    txBox = add_body_textbox(slide, 7.12, 2.25, 5.05, 2.9)
    tf = txBox.text_frame
    tf.word_wrap = True
    achievements = graph_summary.get("achievements", [])
    for item in achievements[:3]:
        set_paragraph(tf, f"- {item}", size=11.5, color=COLOR_TEXT_PRIMARY, spacing_after=6)

    # ══════════════════════════════════════════════════════════════════════════
    #  SLIDE 6: Pass / Fail Overview
    # ══════════════════════════════════════════════════════════════════════════
    pass_fail = ctx.get("pass_fail_chart", {})
    pf_labels = pass_fail.get("labels", [])
    pf_data = pass_fail.get("datasets", [{}])[0].get("data", [])

    slide = new_slide("Pass / Fail Overview", footer_text=filter_str)
    add_subtitle(
        slide,
        f"Threshold: 40% (Pass = ≥40%, Fail = <40%)",
        top=1.0,
        height=0.4,
    )
    pf_header = ["Category", "Count", "Percentage"]
    pf_table_data = [pf_header]
    for i, label in enumerate(pf_labels):
        count = pf_data[i] if i < len(pf_data) else 0
        pct = round((count / total * 100), 1) if total else 0
        pf_table_data.append([label, str(count), f"{pct}%"])

    tbl = add_table(slide, pf_table_data, top=1.6, height=2.0)
    # Colour pass green, fail red
    if len(pf_table_data) > 1:
        for r_idx in range(1, len(pf_table_data)):
            cell = tbl.cell(r_idx, 0)
            bg = COLOR_SUCCESS if "Pass" in pf_table_data[r_idx][0] else COLOR_DANGER
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(bg)
            for para in cell.text_frame.paragraphs:
                para.font.color.rgb = _rgb(COLOR_WHITE)
                para.font.bold = True

    # Add summary text
    txBox = add_body_textbox(slide, 0.5, 4.0, 12, 1.5)
    tf = txBox.text_frame
    tf.word_wrap = True
    pass_val = pf_data[0] if len(pf_data) > 0 else 0
    fail_val = pf_data[1] if len(pf_data) > 1 else 0
    set_paragraph(
        tf,
        f"{pass_val} students passed ({pass_rate}%) and {fail_val} students "
        f"scored below 40%. Focus intervention on the {below_40} result(s) below the threshold.",
        size=14,
        color=COLOR_TEXT_PRIMARY,
    )

    # ══════════════════════════════════════════════════════════════════════════
    #  SLIDE 7: Regression Trend
    # ══════════════════════════════════════════════════════════════════════════
    reg_chart = ctx.get("regression_chart", {})
    slope = reg_chart.get("slope")
    scatter = reg_chart.get("scatter", [])

    slide = new_slide("Performance Trend — Regression Analysis", footer_text=filter_str)
    txBox = add_body_textbox(slide, 0.5, 1.5, 12, 4)
    tf = txBox.text_frame
    tf.word_wrap = True

    run_terminals = ctx.get("run_terminals", [])
    if slope is not None:
        direction = "improving 📈" if slope > 0.5 else ("declining 📉" if slope < -0.5 else "stable ➡️")
        set_paragraph(
            tf,
            f"Trend Direction: {direction}",
            size=20,
            bold=True,
            color=COLOR_PRIMARY,
            spacing_after=12,
        )
        set_paragraph(
            tf,
            f"The regression slope is {slope:.2f}% per terminal.",
            size=16,
            color=COLOR_TEXT_PRIMARY,
            spacing_after=8,
        )
        if slope > 0.5:
            set_paragraph(
                tf,
                "The class is showing improvement over time. Continue current teaching strategies.",
                size=14,
                color=COLOR_SUCCESS,
                spacing_after=4,
            )
        elif slope < -0.5:
            set_paragraph(
                tf,
                "The trend is declining. Early intervention and targeted remediation are recommended.",
                size=14,
                color=COLOR_DANGER,
                spacing_after=4,
            )
        else:
            set_paragraph(
                tf,
                "Performance is steady across terminals. Look for individual student variation.",
                size=14,
                color=COLOR_INFO,
                spacing_after=4,
            )
    else:
        set_paragraph(
            tf,
            "Trend Analysis: Insufficient Data",
            size=20,
            bold=True,
            color=COLOR_TEXT_SECONDARY,
            spacing_after=12,
        )
        set_paragraph(
            tf,
            "At least two terminals of data are required to calculate a regression trend. "
            "Add marks for additional terminals to enable this analysis.",
            size=14,
            color=COLOR_TEXT_SECONDARY,
            spacing_after=4,
        )

    # ══════════════════════════════════════════════════════════════════════════
    #  SLIDE 8: Student Trend Table
    # ══════════════════════════════════════════════════════════════════════════
    trend_rows = ctx.get("trend_rows", [])

    slide = new_slide("Student Trend Table", footer_text=filter_str)
    add_subtitle(
        slide,
        "Per-terminal marks and trend direction for each student",
        top=1.0,
        height=0.4,
    )
    if trend_rows:
        term_labels = run_terminals
        trend_header = ["Student Name"] + list(term_labels) + ["Trend"]
        trend_table_data = [trend_header]
        for row in trend_rows:
            student = row.get("student")
            marks_list = row.get("marks_list", [])
            arrow = row.get("arrow", "same")
            arrow_symbol = {"up": "↑ Improving", "down": "↓ Declining", "same": "→ Stable"}.get(arrow, "→")
            name_str = f"{student.name} ({getattr(student, 'roll_number', '')})" if hasattr(student, 'name') else "Unknown"
            data_row = [name_str]
            for m in marks_list:
                data_row.append(f"{m}%" if m is not None else "—")
            data_row.append(arrow_symbol)
            trend_table_data.append(data_row)

        max_rows = min(len(trend_table_data), 16)  # limit to 16 rows + header
        truncated_data = trend_table_data[:max_rows]
        tbl = add_table(slide, truncated_data, top=1.6, height=min(max_rows * 0.38, 5.0))

        # Colour-code trend column
        if len(truncated_data) > 1:
            for r_idx in range(1, len(truncated_data)):
                cell = tbl.cell(r_idx, len(term_labels) + 1)
                arrow_text = truncated_data[r_idx][-1]
                if "Improving" in arrow_text:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(COLOR_SUCCESS)
                elif "Declining" in arrow_text:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(COLOR_DANGER)
                for para in cell.text_frame.paragraphs:
                    para.font.color.rgb = _rgb(COLOR_WHITE)
                    para.font.bold = True
        if len(trend_table_data) > max_rows:
            txBox = add_body_textbox(slide, 0.5, 6.5, 12, 0.4)
            tf = txBox.text_frame
            set_paragraph(
                tf,
                f"... and {len(trend_table_data) - max_rows} more students. "
                f"Total: {len(trend_rows)} students.",
                size=11,
                color=COLOR_TEXT_MUTED,
            )
    else:
        txBox = add_body_textbox(slide, 0.5, 2.0, 12, 1)
        tf = txBox.text_frame
        set_paragraph(
            tf,
            "No student trend data available for the current filters.",
            size=14,
            color=COLOR_TEXT_SECONDARY,
        )

    # ══════════════════════════════════════════════════════════════════════════
    #  SLIDE 9: Strengths & Weaknesses
    # ══════════════════════════════════════════════════════════════════════════
    # Compute strongest/focus subject from subject_snapshot
    chart_subjects = [
        {"name": s.get("name", "?"), "pct": s.get("pct")}
        for s in ctx.get("subject_snapshot", [])
    ]
    improving_count = ctx.get("improving_count", 0)
    declining_count = ctx.get("declining_count", 0)
    stable_count = ctx.get("stable_count", 0)

    # Recompute strongest/focus from available data
    strongest = None
    focus = None
    if chart_subjects:
        valid = [s for s in chart_subjects if s["pct"] is not None]
        if valid:
            strongest = max(valid, key=lambda x: x["pct"])
            focus = min(valid, key=lambda x: x["pct"])

    slide = new_slide("Strengths & Weaknesses", footer_text=filter_str)
    sw_data = [
        ["Category", "Detail"],
        [
            "Strongest Subject",
            f"{strongest['name']} ({strongest['pct']}%)" if strongest else "N/A",
        ],
        [
            "Area for Focus",
            f"{focus['name']} ({focus['pct']}%)" if focus and focus != strongest else "N/A",
        ],
        ["Students Improving", str(improving_count)],
        ["Students Declining", str(declining_count)],
        ["Students Stable", str(stable_count)],
    ]
    tbl = add_table(slide, sw_data, top=1.6, height=3.0)

    # Colour-code the detail cells
    if strongest:
        cell = tbl.cell(1, 1)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(COLOR_SUCCESS)
        for para in cell.text_frame.paragraphs:
            para.font.color.rgb = _rgb(COLOR_WHITE)
            para.font.bold = True
    if focus and focus != strongest:
        cell = tbl.cell(2, 1)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(COLOR_DANGER)
        for para in cell.text_frame.paragraphs:
            para.font.color.rgb = _rgb(COLOR_WHITE)
            para.font.bold = True
    # Improving/declining cells
    cell_imp = tbl.cell(3, 1)
    cell_imp.fill.solid()
    cell_imp.fill.fore_color.rgb = _rgb(COLOR_SUCCESS)
    for para in cell_imp.text_frame.paragraphs:
        para.font.color.rgb = _rgb(COLOR_WHITE)
        para.font.bold = True
    cell_dec = tbl.cell(4, 1)
    cell_dec.fill.solid()
    cell_dec.fill.fore_color.rgb = _rgb(COLOR_DANGER if declining_count > 0 else COLOR_TEXT_MUTED)
    for para in cell_dec.text_frame.paragraphs:
        para.font.color.rgb = _rgb(COLOR_WHITE)
        para.font.bold = True

    add_panel(slide, 0.5, 4.85, 12, 1.45)
    add_section_label(slide, "Recommended Next Step", 0.7, 5.04, color=COLOR_WARNING)
    txBox = add_body_textbox(slide, 0.75, 5.42, 11.4, 0.62)
    tf = txBox.text_frame
    tf.word_wrap = True
    recommendation = (graph_summary.get("recommendations") or ["Use the weakest area as the next short reteaching target."])[0]
    set_paragraph(tf, recommendation, size=12, color=COLOR_TEXT_PRIMARY)

    # ══════════════════════════════════════════════════════════════════════════
    #  SLIDE 10: Insights & Recommendations
    # ══════════════════════════════════════════════════════════════════════════
    slide = new_slide("Insights & Recommendations", footer_text=filter_str)
    txBox = add_body_textbox(slide, 0.5, 1.5, 12, 5.0)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Headline
    headline = graph_summary.get("headline", "Analytics summary available on screen.")
    set_paragraph(tf, "📊 Overview", size=18, bold=True, color=COLOR_PRIMARY, spacing_after=8)
    set_paragraph(tf, headline, size=14, color=COLOR_TEXT_PRIMARY, spacing_after=12)

    # Graph Reading
    graph_reading = graph_summary.get("graph_reading", [])
    if graph_reading:
        set_paragraph(tf, "📈 How to Read the Charts", size=18, bold=True, color=COLOR_INFO, spacing_after=8)
        for item in graph_reading:
            set_paragraph(tf, f"• {item}", size=13, color=COLOR_TEXT_PRIMARY, spacing_after=4)
        set_paragraph(tf, "", size=6)

    # Achievements
    achievements = graph_summary.get("achievements", [])
    if achievements:
        set_paragraph(tf, "🏆 Achievements & Milestones", size=18, bold=True, color=COLOR_SUCCESS, spacing_after=8)
        for item in achievements:
            set_paragraph(tf, f"• {item}", size=13, color=COLOR_TEXT_PRIMARY, spacing_after=4)
        set_paragraph(tf, "", size=6)

    # Recommendations
    recommendations = graph_summary.get("recommendations", [])
    if recommendations:
        set_paragraph(tf, "💡 Teaching Recommendations", size=18, bold=True, color=COLOR_EXCELLENCE, spacing_after=8)
        for item in recommendations:
            set_paragraph(tf, f"• {item}", size=13, color=COLOR_TEXT_PRIMARY, spacing_after=4)

    # ══════════════════════════════════════════════════════════════════════════
    #  SLIDE 11: Detailed Insight Cards (conditional)
    # ══════════════════════════════════════════════════════════════════════════
    enriched = ctx.get("teacher_enriched_insights", {})
    if enriched:
        slide = new_slide("Detailed Insights", footer_text=filter_str)
        txBox = add_body_textbox(slide, 0.5, 1.5, 12, 5.0)
        tf = txBox.text_frame
        tf.word_wrap = True

        sections = [
            ("strength", "💪 Strengths", COLOR_SUCCESS),
            ("weakness", "⚠️ Weaknesses", COLOR_DANGER),
            ("trend", "📊 Trends", COLOR_INFO),
            ("compare", "🔍 Comparisons", COLOR_WARNING),
            ("recommend", "💡 Recommendations", COLOR_EXCELLENCE),
            ("milestone", "🎯 Milestones", "#06b6d4"),
        ]

        first = True
        for key, label, color in sections:
            items = enriched.get(key, [])
            if items:
                if not first:
                    set_paragraph(tf, "", size=4)
                set_paragraph(tf, label, size=16, bold=True, color=color, spacing_after=6)
                for item in items[:5]:  # max 5 per section
                    title = item.get("title", "")
                    description = item.get("description", "")
                    if title:
                        set_paragraph(tf, f"▸ {title}", size=13, bold=True, color=COLOR_TEXT_PRIMARY, spacing_after=2)
                    if description:
                        set_paragraph(tf, f"  {description}", size=12, color=COLOR_TEXT_SECONDARY, spacing_after=4)
                first = False

    # ══════════════════════════════════════════════════════════════════════════
    #  SLIDE 12: Metadata / Footer
    # ══════════════════════════════════════════════════════════════════════════
    slide = new_slide(
        "Report Metadata",
        subtitle_text="Complete filter parameters and generation information",
        footer_text=f"Generated by AcadStat — {today_str}",
    )
    txBox = add_body_textbox(slide, 0.5, 2.0, 12, 4.5)
    tf = txBox.text_frame
    tf.word_wrap = True
    meta_items = [
        ("Generated", today_str),
        ("Teacher", teacher_name),
        ("Analysis Mode", mode_display),
        ("Subject(s)", subj_display),
        ("Terminal(s)", term_display),
        ("Class / Section", cs_display),
    ]
    if selected_student:
        meta_items.append(("Student", student_display))
    meta_items += [
        ("Total Results Analysed", str(total)),
        ("Pass Rate", f"{pass_rate}%"),
        ("Average Score", f"{avg_pct}%"),
    ]
    for label, value in meta_items:
        set_paragraph(tf, f"{label}: {value}", size=14, color=COLOR_TEXT_PRIMARY, spacing_after=4)

    set_paragraph(tf, "", size=8)
    set_paragraph(
        tf,
        "Disclaimer: This report was generated automatically by AcadStat Academic "
        "Management System. Data reflects marks entered by teachers for students "
        "assigned to them. Always verify against original records.",
        size=11,
        color=COLOR_TEXT_MUTED,
    )

    # ── Save and return ─────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    resp = HttpResponse(
        buf.read(),
        content_type=(
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation"
        ),
    )
    resp["Content-Disposition"] = f'attachment; filename="{_filename(ctx, "pptx")}"'
    return resp


# ═════════════════════════════════════════════════════════════════════════════
#  PDF EXPORT
# ═════════════════════════════════════════════════════════════════════════════

def export_subject_analytics_pdf(request, context):
    """
    Build a PDF report from the subject analytics context using ReportLab.
    Returns an HttpResponse with content-type 'application/pdf'.
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm, mm
        from reportlab.lib import colors as rl_colors
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import (
            SimpleDocTemplate,
            Table,
            TableStyle,
            Paragraph,
            Spacer,
            PageBreak,
            HRFlowable,
        )
        from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
    except ImportError:
        logger.exception("reportlab is not installed")
        return HttpResponse(
            "PDF generation requires reportlab.",
            status=500,
        )

    ctx = context
    teacher_name = _teacher_display_name(getattr(request, "teacher", None))
    filter_str = _filter_label(ctx)
    filter_rows = _filter_rows(ctx)
    today_str = date.today().strftime("%B %d, %Y")
    total = ctx.get("total", 0)
    avg_pct = ctx.get("avg_pct", 0)
    highest_pct = ctx.get("highest_pct", 0)
    lowest_pct = ctx.get("lowest_pct", 0)
    pass_rate = ctx.get("pass_rate", 0)
    pass_count = ctx.get("pass_count", 0)
    below_40 = ctx.get("below_40", 0)
    mode_display = "Single Analysis" if ctx.get("analysis_mode") == "single" else "Bulk Analysis"
    subj_obj = ctx.get("subject_obj")
    subj_display = subj_obj.name if subj_obj else "All subjects"
    term_display = ctx.get("terminal", "all")
    cs_display = ctx.get("class_section", "").replace("|", " - ") or "All sections"

    # ── Styles ────────────────────────────────────────────────────────
    styles = getSampleStyleSheet()
    accent = rl_colors.HexColor(COLOR_PRIMARY)
    grey = rl_colors.HexColor(COLOR_TEXT_SECONDARY)
    muted = rl_colors.HexColor(COLOR_TEXT_MUTED)

    def RLColor(hex_color):
        return rl_colors.HexColor(hex_color)

    def P(text, style_name="Normal", size=10, color=None, bold=False, align=TA_LEFT):
        color = color or rl_colors.HexColor(COLOR_TEXT_PRIMARY)
        st = ParagraphStyle(
            "_tmp",
            parent=styles.get(style_name, styles["Normal"]),
            fontSize=size,
            textColor=color,
            alignment=align,
            spaceAfter=6,
            leading=size * 1.4,
        )
        if bold:
            st.fontName = "Helvetica-Bold"
        return Paragraph(text, st)

    def make_table(data, col_widths=None, header_bg=COLOR_PRIMARY):
        """Create a styled table. data[0] = header."""
        tbl = Table(data, colWidths=col_widths, repeatRows=1)
        style_cmds = [
            ("BACKGROUND", (0, 0), (-1, 0), RLColor(header_bg)),
            ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("GRID", (0, 0), (-1, -1), 0.5, rl_colors.grey),
            ("ALIGN", (0, 0), (-1, -1), "CENTER"),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ]
        # Alternating row colours
        for i in range(1, len(data)):
            if i % 2 == 0:
                style_cmds.append(
                    ("BACKGROUND", (0, i), (-1, i), RLColor("#F8FAFC"))
                )
        tbl.setStyle(TableStyle(style_cmds))
        return tbl

    def section_title(text, subtitle=None):
        elements.append(Spacer(1, 0.28 * cm))
        elements.append(HRFlowable(width="100%", thickness=1.2, color=RLColor(COLOR_HEADER_BG), spaceAfter=0.16 * cm))
        elements.append(P(text, "Normal", size=14, color=COLOR_PRIMARY, bold=True))
        if subtitle:
            elements.append(P(subtitle, "Normal", size=9, color=COLOR_TEXT_SECONDARY))
        elements.append(Spacer(1, 0.12 * cm))

    def metric_card_table(cards):
        data = [[P(label, size=8, color=COLOR_TEXT_SECONDARY, bold=True, align=TA_CENTER) for label, _, _ in cards],
                [P(str(value), size=16, color=color, bold=True, align=TA_CENTER) for _, value, color in cards]]
        tbl = Table(data, colWidths=[3.35 * cm] * len(cards))
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), RLColor("#FFFFFF")),
            ("BOX", (0, 0), (-1, -1), 0.7, RLColor(COLOR_HEADER_BG)),
            ("INNERGRID", (0, 0), (-1, -1), 0.5, RLColor(COLOR_HEADER_BG)),
            ("TOPPADDING", (0, 0), (-1, -1), 7),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ]))
        return tbl

    def page_canvas(canvas, doc_obj):
        canvas.saveState()
        page_w, page_h = A4
        canvas.setFillColor(RLColor(COLOR_PRIMARY))
        canvas.rect(0, page_h - 7 * mm, page_w, 7 * mm, fill=1, stroke=0)
        canvas.setFillColor(RLColor(COLOR_SUCCESS))
        canvas.rect(0, page_h - 9 * mm, page_w, 2 * mm, fill=1, stroke=0)
        canvas.setFont("Helvetica-Bold", 8)
        canvas.setFillColor(RLColor(COLOR_TEXT_PRIMARY))
        canvas.drawString(12 * mm, 7 * mm, f"AcadStat Subject Analytics | {teacher_name}")
        canvas.setFont("Helvetica", 7)
        canvas.setFillColor(RLColor(COLOR_TEXT_MUTED))
        footer = filter_str
        if len(footer) > 120:
            footer = footer[:117] + "..."
        canvas.drawCentredString(page_w / 2, 7 * mm, footer)
        canvas.drawRightString(page_w - 12 * mm, 7 * mm, f"Page {doc_obj.page}")
        canvas.restoreState()

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        topMargin=1.25 * cm,
        bottomMargin=1.25 * cm,
        leftMargin=1.15 * cm,
        rightMargin=1.15 * cm,
        title="Subject Analytics Report",
    )

    elements = []

    # ── Page 1: Title Page ─────────────────────────────────────────────
    elements.append(Spacer(1, 0.35 * cm))
    elements.append(
        P("Subject Analytics Report", "Title", size=24, color=COLOR_PRIMARY, bold=True, align=TA_CENTER)
    )
    elements.append(Spacer(1, 0.2 * cm))
    elements.append(
        P(f"Teacher: {teacher_name}", "Normal", size=14, color=COLOR_TEXT_SECONDARY, align=TA_CENTER)
    )
    elements.append(
        P(f"Generated: {today_str}", "Normal", size=12, color=COLOR_TEXT_SECONDARY, align=TA_CENTER)
    )
    elements.append(Spacer(1, 0.22 * cm))
    cover_filters = [["Applied Filter", "Selected Value"]] + [[label, value] for label, value in filter_rows]
    elements.append(make_table(cover_filters, col_widths=[5.2 * cm, 10.2 * cm]))
    elements.append(Spacer(1, 0.28 * cm))
    elements.append(HRFlowable(width="100%", thickness=2, color=accent, spaceBefore=0.1 * cm, spaceAfter=0.25 * cm))
    elements.append(
        P(
            "This report provides a complete summary of student performance based "
            "on the currently applied filters. It includes KPIs, score distribution, "
            "terminal performance, subject breakdowns, trend analysis, and teaching insights.",
            "Normal",
            size=11,
            color=COLOR_TEXT_SECONDARY,
            align=TA_CENTER,
        )
    )
    elements.append(Spacer(1, 0.3 * cm))

    # If no data
    if total == 0:
        elements.append(
            P("No Data Available", "Normal", size=18, bold=True, color=COLOR_DANGER, align=TA_CENTER)
        )
        elements.append(
            P(
                "No marks data was found for the current filter combination. "
                "Try adjusting your filters or add marks from the 'Add Marks' page.",
                "Normal",
                size=12,
                color=COLOR_TEXT_SECONDARY,
                align=TA_CENTER,
            )
        )
        doc.build(elements, onFirstPage=page_canvas, onLaterPages=page_canvas)
        buf.seek(0)
        resp = HttpResponse(buf.read(), content_type="application/pdf")
        resp["Content-Disposition"] = f'attachment; filename="{_filename(ctx, "pdf")}"'
        return resp

    elements.append(Spacer(1, 0.35 * cm))

    # ── Page 2: KPI Summary ────────────────────────────────────────────
    elements.append(P("Key Performance Indicators", "Normal", size=18, color=COLOR_PRIMARY, bold=True))
    elements.append(Spacer(1, 0.3 * cm))
    kpi_data = [
        ["Metric", "Value"],
        ["Average Score", f"{avg_pct}%"],
        ["Highest Score", f"{highest_pct}%"],
        ["Lowest Score", f"{lowest_pct}%"],
        ["Pass Rate", f"{pass_rate}% ({pass_count}/{total})"],
        ["Below 40%", str(below_40)],
        ["Analysis Mode", "Single Analysis" if ctx.get("analysis_mode") == "single" else "Bulk Analysis"],
    ]
    if ctx.get("selected_student"):
        kpi_data.append(["Student", ctx["selected_student"].name])
    elements.append(make_table(kpi_data, col_widths=[6 * cm, 6 * cm]))
    elements.append(Spacer(1, 0.5 * cm))
    elements.append(HRFlowable(width="100%", thickness=1, color=RLColor(COLOR_HEADER_BG), spaceAfter=0.3 * cm))

    # ── Page 3: Score Distribution ─────────────────────────────────────
    elements.append(Spacer(1, 0.25 * cm))
    elements.append(P("Score Distribution", "Normal", size=18, color=COLOR_PRIMARY, bold=True))
    elements.append(
        P("Distribution of student scores across grade bands", "Normal", size=11, color=COLOR_TEXT_SECONDARY)
    )
    elements.append(Spacer(1, 0.3 * cm))
    dist_chart = ctx.get("distribution_chart", {})
    dist_labels = dist_chart.get("labels", [])
    dist_data = dist_chart.get("datasets", [{}])[0].get("data", [])
    dist_bg = dist_chart.get("datasets", [{}])[0].get("backgroundColor", [])
    dist_table_data = [["Grade Band", "Count", "Percentage"]]
    for i, label in enumerate(dist_labels):
        count = dist_data[i] if i < len(dist_data) else 0
        pct_str = f"{round(count / total * 100, 1)}%" if total else "0%"
        dist_table_data.append([label, str(count), pct_str])
    tbl = make_table(dist_table_data, col_widths=[6 * cm, 3 * cm, 3 * cm])
    # Colour first column
    for i in range(1, len(dist_table_data)):
        if i - 1 < len(dist_bg):
            tbl.setStyle(TableStyle([
                ("BACKGROUND", (0, i), (0, i), RLColor(dist_bg[i - 1])),
                ("TEXTCOLOR", (0, i), (0, i), rl_colors.white),
                ("FONTNAME", (0, i), (0, i), "Helvetica-Bold"),
            ]))
    elements.append(tbl)
    elements.append(Spacer(1, 0.3 * cm))
    elements.append(
        P(
            f"• {dist_table_data[1][1] if len(dist_table_data) > 1 else 0} students scored below 40% (fail). "
            f"• {dist_table_data[4][1] if len(dist_table_data) > 4 else 0} students achieved 80%+ (excellent).",
            "Normal",
            size=10,
            color=COLOR_TEXT_PRIMARY,
        )
    )

    # ── Page 4: Terminal Performance ───────────────────────────────────
    elements.append(Spacer(1, 0.25 * cm))
    elements.append(P("Terminal-Wise Performance", "Normal", size=18, color=COLOR_PRIMARY, bold=True))
    elements.append(
        P("Average percentage per terminal by subject", "Normal", size=11, color=COLOR_TEXT_SECONDARY)
    )
    elements.append(Spacer(1, 0.3 * cm))
    terminal_wise = ctx.get("terminal_wise_chart", {})
    term_labels = terminal_wise.get("labels", [])
    term_datasets = terminal_wise.get("datasets", [])
    if term_datasets:
        term_table_data = [["Subject"] + list(term_labels)]
        for ds in term_datasets:
            row = [ds.get("label", "?")]
            for v in ds.get("data", []):
                row.append(f"{v}%" if v is not None else "—")
            term_table_data.append(row)
        tbl = make_table(term_table_data, col_widths=[4 * cm] + [3 * cm] * len(term_labels))
        # Colour percentage cells
        for r_idx in range(1, len(term_table_data)):
            for c_idx in range(1, len(term_table_data[r_idx])):
                cell_val = term_datasets[r_idx - 1].get("data", [])[c_idx - 1] if c_idx - 1 < len(term_datasets[r_idx - 1].get("data", [])) else None
                if cell_val is not None:
                    bg = _pct_color(cell_val)
                    tbl.setStyle(TableStyle([
                        ("BACKGROUND", (c_idx, r_idx), (c_idx, r_idx), RLColor(bg)),
                        ("TEXTCOLOR", (c_idx, r_idx), (c_idx, r_idx), rl_colors.white),
                        ("FONTNAME", (c_idx, r_idx), (c_idx, r_idx), "Helvetica-Bold"),
                    ]))
        elements.append(tbl)
    else:
        elements.append(P("No terminal performance data available.", "Normal", size=11, color=COLOR_TEXT_MUTED))

    # ── Page 5: Subject Snapshot ───────────────────────────────────────
    elements.append(Spacer(1, 0.25 * cm))
    elements.append(P("Subject Snapshot", "Normal", size=18, color=COLOR_PRIMARY, bold=True))
    elements.append(
        P("Subject-wise performance in the latest/selected terminal", "Normal", size=11, color=COLOR_TEXT_SECONDARY)
    )
    elements.append(Spacer(1, 0.3 * cm))
    subj_breakdown = ctx.get("subject_breakdown_chart", {})
    subj_labels = subj_breakdown.get("labels", [])
    subj_values = subj_breakdown.get("datasets", [{}])[0].get("data", [])
    if subj_labels:
        subj_table_data = [["Subject", "Average %"]]
        for i, label in enumerate(subj_labels):
            val = subj_values[i] if i < len(subj_values) else None
            subj_table_data.append([label, f"{val}%" if val is not None else "—"])
        tbl = make_table(subj_table_data, col_widths=[8 * cm, 4 * cm])
        # Colour cells
        for i in range(1, len(subj_table_data)):
            val = subj_values[i - 1] if i - 1 < len(subj_values) else None
            if val is not None:
                bg = _pct_color(val)
                tbl.setStyle(TableStyle([
                    ("BACKGROUND", (1, i), (1, i), RLColor(bg)),
                    ("TEXTCOLOR", (1, i), (1, i), rl_colors.white),
                    ("FONTNAME", (1, i), (1, i), "Helvetica-Bold"),
                ]))
        elements.append(tbl)
    else:
        elements.append(P("No subject data available.", "Normal", size=11, color=COLOR_TEXT_MUTED))

    # ── Page 6: Pass/Fail & Regression ─────────────────────────────────
    elements.append(Spacer(1, 0.25 * cm))
    elements.append(P("Pass / Fail & Trend Analysis", "Normal", size=18, color=COLOR_PRIMARY, bold=True))
    elements.append(Spacer(1, 0.3 * cm))
    pf_chart = ctx.get("pass_fail_chart", {})
    pf_labels = pf_chart.get("labels", [])
    pf_data = pf_chart.get("datasets", [{}])[0].get("data", [])
    if pf_labels:
        pf_table_data = [["Category", "Count", "Percentage"]]
        for i, label in enumerate(pf_labels):
            count = pf_data[i] if i < len(pf_data) else 0
            pct_str = f"{round(count / total * 100, 1)}%" if total else "0%"
            pf_table_data.append([label, str(count), pct_str])
        tbl = make_table(pf_table_data, col_widths=[5 * cm, 3 * cm, 4 * cm])
        # Colour pass green, fail red
        for i in range(1, len(pf_table_data)):
            bg = COLOR_SUCCESS if "Pass" in pf_table_data[i][0] else COLOR_DANGER
            tbl.setStyle(TableStyle([
                ("BACKGROUND", (0, i), (0, i), RLColor(bg)),
                ("TEXTCOLOR", (0, i), (0, i), rl_colors.white),
                ("FONTNAME", (0, i), (0, i), "Helvetica-Bold"),
            ]))
        elements.append(tbl)
    elements.append(Spacer(1, 0.3 * cm))

    # Regression
    reg_chart = ctx.get("regression_chart", {})
    slope = reg_chart.get("slope")
    elements.append(
        P("Performance Trend", "Normal", size=14, color=COLOR_PRIMARY, bold=True)
    )
    if slope is not None:
        direction = "improving" if slope > 0.5 else ("declining" if slope < -0.5 else "stable")
        trend_color = COLOR_SUCCESS if slope > 0.5 else (COLOR_DANGER if slope < -0.5 else COLOR_INFO)
        elements.append(
            P(
                f"Trend: {direction} ({slope:.2f}% per terminal)",
                "Normal",
                size=12,
                color=trend_color,
                bold=True,
            )
        )
        if slope > 0.5:
            elements.append(P("The class is showing improvement over time. Continue current strategies.", "Normal", size=10, color=COLOR_TEXT_PRIMARY))
        elif slope < -0.5:
            elements.append(P("The trend is declining. Early intervention and targeted remediation are recommended.", "Normal", size=10, color=COLOR_DANGER))
        else:
            elements.append(P("Performance is steady across terminals.", "Normal", size=10, color=COLOR_TEXT_PRIMARY))
    else:
        elements.append(
            P("Insufficient data to calculate trend (need at least 2 terminals).", "Normal", size=10, color=COLOR_TEXT_MUTED)
        )

    # ── Page 7: Student Trend Table ────────────────────────────────────
    trend_rows = ctx.get("trend_rows", [])
    run_terminals = ctx.get("run_terminals", [])
    elements.append(Spacer(1, 0.25 * cm))
    elements.append(P("Student Trend Table", "Normal", size=18, color=COLOR_PRIMARY, bold=True))
    elements.append(
        P("Per-terminal marks and trend direction", "Normal", size=11, color=COLOR_TEXT_SECONDARY)
    )
    elements.append(Spacer(1, 0.3 * cm))
    if trend_rows:
        trend_table_data = [["Student"] + list(run_terminals) + ["Trend"]]
        for row in trend_rows[:25]:  # max 25 students in PDF
            student = row.get("student")
            name_str = student.name if hasattr(student, "name") else "Unknown"
            data_row = [name_str]
            for m in row.get("marks_list", []):
                data_row.append(f"{m}%" if m is not None else "—")
            arrow = row.get("arrow", "same")
            arrow_str = {"up": "↑", "down": "↓", "same": "→"}.get(arrow, "→")
            data_row.append(arrow_str)
            trend_table_data.append(data_row)

        col_w = [4 * cm] + [2 * cm] * len(run_terminals) + [2 * cm]
        tbl = make_table(trend_table_data, col_widths=col_w)
        # Colour trend column
        for i in range(1, len(trend_table_data)):
            trend_text = trend_table_data[i][-1]
            if "↑" in trend_text:
                tbl.setStyle(TableStyle([
                    ("BACKGROUND", (-1, i), (-1, i), RLColor(COLOR_SUCCESS)),
                    ("TEXTCOLOR", (-1, i), (-1, i), rl_colors.white),
                ]))
            elif "↓" in trend_text:
                tbl.setStyle(TableStyle([
                    ("BACKGROUND", (-1, i), (-1, i), RLColor(COLOR_DANGER)),
                    ("TEXTCOLOR", (-1, i), (-1, i), rl_colors.white),
                ]))
        elements.append(tbl)
        if len(trend_rows) > 25:
            elements.append(
                P(
                    f"Showing 25 of {len(trend_rows)} students. Full data available on screen.",
                    "Normal",
                    size=9,
                    color=COLOR_TEXT_MUTED,
                )
            )
    else:
        elements.append(P("No student trend data available.", "Normal", size=11, color=COLOR_TEXT_MUTED))

    # ── Page 8: Strengths & Weaknesses ─────────────────────────────────
    elements.append(Spacer(1, 0.25 * cm))
    elements.append(P("Strengths & Weaknesses", "Normal", size=18, color=COLOR_PRIMARY, bold=True))
    elements.append(Spacer(1, 0.3 * cm))
    chart_subjects_list = ctx.get("subject_snapshot", [])
    strongest = max(chart_subjects_list, key=lambda x: x["pct"]) if chart_subjects_list else None
    focus = min(chart_subjects_list, key=lambda x: x["pct"]) if chart_subjects_list else None
    improving_count = ctx.get("improving_count", 0)
    declining_count = ctx.get("declining_count", 0)
    stable_count = ctx.get("stable_count", 0)

    sw_data = [
        ["Category", "Detail"],
        ["Strongest Subject", f"{strongest['name']} ({strongest['pct']}%)" if strongest else "N/A"],
        ["Area for Focus", f"{focus['name']} ({focus['pct']}%)" if focus and focus != strongest else "N/A"],
        ["Students Improving", str(improving_count)],
        ["Students Declining", str(declining_count)],
        ["Students Stable", str(stable_count)],
    ]
    tbl = make_table(sw_data, col_widths=[6 * cm, 6 * cm])
    if strongest:
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (1, 1), (1, 1), RLColor(COLOR_SUCCESS)),
            ("TEXTCOLOR", (1, 1), (1, 1), rl_colors.white),
            ("FONTNAME", (1, 1), (1, 1), "Helvetica-Bold"),
        ]))
    if focus and focus != strongest:
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (1, 2), (1, 2), RLColor(COLOR_DANGER)),
            ("TEXTCOLOR", (1, 2), (1, 2), rl_colors.white),
            ("FONTNAME", (1, 2), (1, 2), "Helvetica-Bold"),
        ]))
    elements.append(tbl)

    # ── Page 9: Teacher Insights ───────────────────────────────────────
    elements.append(Spacer(1, 0.25 * cm))
    elements.append(P("Teacher Insights & Recommendations", "Normal", size=18, color=COLOR_PRIMARY, bold=True))
    elements.append(Spacer(1, 0.3 * cm))
    graph_summary = ctx.get("teacher_graph_summary", {})

    headline = graph_summary.get("headline", "")
    if headline:
        elements.append(P(headline, "Normal", size=12, color=COLOR_TEXT_PRIMARY))
        elements.append(Spacer(1, 0.2 * cm))

    for section_title, key in [
        ("How to Read the Charts", "graph_reading"),
        ("Achievements & Milestones", "achievements"),
        ("Teaching Recommendations", "recommendations"),
    ]:
        items = graph_summary.get(key, [])
        if items:
            elements.append(P(section_title, "Normal", size=14, color=COLOR_PRIMARY, bold=True))
            for item in items:
                elements.append(P(f"• {item}", "Normal", size=10, color=COLOR_TEXT_PRIMARY))
            elements.append(Spacer(1, 0.2 * cm))

    # ── Page 10: Appendix ──────────────────────────────────────────────
    elements.append(Spacer(1, 0.25 * cm))
    elements.append(P("Appendix — Report Metadata", "Normal", size=18, color=COLOR_PRIMARY, bold=True))
    elements.append(Spacer(1, 0.3 * cm))
    meta_data = [
        ["Parameter", "Value"],
        ["Generation Date", today_str],
        ["Teacher", teacher_name],
        ["Analysis Mode", mode_display := "Single Analysis" if ctx.get("analysis_mode") == "single" else "Bulk Analysis"],
        ["Subject(s)", subj_display := (ctx.get("subject_obj").name if ctx.get("subject_obj") else "All")],
        ["Terminal(s)", ctx.get("terminal", "all")],
        ["Class / Section", ctx.get("class_section", "").replace("|", " — ") or "All"],
    ]
    if ctx.get("selected_student"):
        meta_data.append(["Student", ctx["selected_student"].name])
    meta_data += [
        ["Results Analysed", str(total)],
        ["Pass Rate", f"{pass_rate}%"],
        ["Average Score", f"{avg_pct}%"],
    ]
    elements.append(make_table(meta_data, col_widths=[5 * cm, 7 * cm]))
    elements.append(Spacer(1, 0.5 * cm))
    elements.append(HRFlowable(width="100%", thickness=1, color=RLColor(COLOR_HEADER_BG), spaceBefore=0.3 * cm, spaceAfter=0.3 * cm))
    elements.append(
        P(
            "Disclaimer: This report was generated automatically by AcadStat Academic "
            "Management System. Data reflects marks entered by teachers for students "
            "assigned to them. Always verify against original records.",
            "Normal",
            size=8,
            color=COLOR_TEXT_MUTED,
        )
    )

    # ── Build ──────────────────────────────────────────────────────────
    doc.build(elements, onFirstPage=page_canvas, onLaterPages=page_canvas)
    buf.seek(0)

    resp = HttpResponse(buf.read(), content_type="application/pdf")
    resp["Content-Disposition"] = f'attachment; filename="{_filename(ctx, "pdf")}"'
    return resp
